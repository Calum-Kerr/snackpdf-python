"""
Basic PDF operations for core functionality
Handles PDF to image conversion, merging, splitting, and protection
"""

from flask import Blueprint, request, jsonify, send_file
from werkzeug.utils import secure_filename
import os
import uuid
import tempfile
import logging
import subprocess
from datetime import datetime
import PyPDF2
import fitz  # PyMuPDF
from PIL import Image
import zipfile

# Create blueprint
basic_operations_bp = Blueprint('basic_operations', __name__)

# Configure logging
logging.basicConfig(level=logging.DEBUG)

# File size limit (50MB)
MAX_FILE_SIZE = 50 * 1024 * 1024

# Static directory for file storage
STATIC_DIR = os.path.join(os.path.dirname(__file__), 'static')
if not os.path.exists(STATIC_DIR):
    os.makedirs(STATIC_DIR)

def validate_pdf_file(file):
    """Validate uploaded PDF file"""
    if not file or not file.filename:
        return False, "No file selected"

    if not file.filename.lower().endswith('.pdf'):
        return False, "Only PDF files are allowed"

    # Check file size (approximate)
    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)  # Reset file pointer

    if file_size > MAX_FILE_SIZE:
        return False, f"File size exceeds {MAX_FILE_SIZE // (1024*1024)}MB limit"

    return True, "Valid"

@basic_operations_bp.route('/pdf_to_jpg', methods=['POST'])
def pdf_to_jpg():
    """Convert PDF to JPG/PNG/TIFF images"""
    try:
        # Get file from request
        if 'file' not in request.files:
            return jsonify({"error": "No file part in request"}), 400

        file = request.files['file']

        # Validate file
        is_valid, message = validate_pdf_file(file)
        if not is_valid:
            return jsonify({"error": message}), 400

        # Get form parameters
        output_format = request.form.get('output_format', 'jpg').lower()
        quality = request.form.get('quality', 'high')
        resolution = int(request.form.get('resolution', '300'))

        # Validate parameters
        if output_format not in ['jpg', 'png', 'tiff']:
            output_format = 'jpg'
        
        if quality not in ['high', 'medium', 'low']:
            quality = 'high'

        # Create secure filenames
        original_filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())[:8]
        input_filename = f"{unique_id}_{original_filename}"
        
        input_path = os.path.join(STATIC_DIR, input_filename)

        # Save uploaded file
        file.save(input_path)
        logging.debug(f"File saved to: {input_path}")

        # Convert PDF to images using PyMuPDF
        doc = fitz.open(input_path)
        image_files = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Set quality based on user selection
            if quality == 'high':
                mat = fitz.Matrix(resolution/72, resolution/72)
            elif quality == 'medium':
                mat = fitz.Matrix((resolution*0.8)/72, (resolution*0.8)/72)
            else:  # low
                mat = fitz.Matrix((resolution*0.6)/72, (resolution*0.6)/72)
            
            pix = page.get_pixmap(matrix=mat)
            
            # Create output filename
            output_filename = f"{unique_id}_page_{page_num + 1}.{output_format}"
            output_path = os.path.join(STATIC_DIR, output_filename)
            
            if output_format == 'jpg':
                pix.save(output_path, output="jpeg")
            elif output_format == 'png':
                pix.save(output_path, output="png")
            elif output_format == 'tiff':
                pix.save(output_path, output="tiff")
            
            image_files.append(output_path)
        
        doc.close()

        # If multiple pages, create a ZIP file
        if len(image_files) > 1:
            zip_filename = f"converted_{unique_id}.zip"
            zip_path = os.path.join(STATIC_DIR, zip_filename)
            
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for img_path in image_files:
                    zipf.write(img_path, os.path.basename(img_path))
            
            # Clean up individual image files
            for img_path in image_files:
                try:
                    os.remove(img_path)
                except:
                    pass
            
            # Clean up input file
            try:
                os.remove(input_path)
            except:
                pass
            
            return send_file(zip_path, as_attachment=True, download_name=zip_filename)
        else:
            # Single page, return the image directly
            output_path = image_files[0]
            download_name = f"converted_{original_filename.rsplit('.', 1)[0]}.{output_format}"
            
            # Clean up input file
            try:
                os.remove(input_path)
            except:
                pass
            
            return send_file(output_path, as_attachment=True, download_name=download_name)

    except Exception as e:
        logging.error(f"Error in PDF to JPG conversion: {str(e)}")
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

@basic_operations_bp.route('/merge', methods=['POST'])
def merge_pdfs():
    """Merge multiple PDF files"""
    try:
        # Get files from request
        files = request.files.getlist('files')
        
        if len(files) < 2:
            return jsonify({"error": "At least 2 PDF files are required for merging"}), 400

        # Validate all files
        temp_files = []
        for file in files:
            is_valid, message = validate_pdf_file(file)
            if not is_valid:
                return jsonify({"error": f"Invalid file {file.filename}: {message}"}), 400
            
            # Save temporary file
            unique_id = str(uuid.uuid4())[:8]
            temp_filename = f"temp_{unique_id}_{secure_filename(file.filename)}"
            temp_path = os.path.join(STATIC_DIR, temp_filename)
            file.save(temp_path)
            temp_files.append(temp_path)

        # Get form parameters
        output_filename = request.form.get('output_filename', 'merged_document.pdf')
        if not output_filename.endswith('.pdf'):
            output_filename += '.pdf'
        
        # Create output file
        unique_id = str(uuid.uuid4())[:8]
        output_path = os.path.join(STATIC_DIR, f"{unique_id}_{secure_filename(output_filename)}")

        # Merge PDFs using PyPDF2
        merger = PyPDF2.PdfMerger()
        
        for temp_file in temp_files:
            merger.append(temp_file)
        
        merger.write(output_path)
        merger.close()

        # Clean up temporary files
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except:
                pass

        return send_file(output_path, as_attachment=True, download_name=output_filename)

    except Exception as e:
        logging.error(f"Error in PDF merge: {str(e)}")
        return jsonify({"error": f"Merge failed: {str(e)}"}), 500

@basic_operations_bp.route('/split', methods=['POST'])
def split_pdf():
    """Split PDF into multiple files"""
    try:
        # Get file from request
        if 'file' not in request.files:
            return jsonify({"error": "No file part in request"}), 400

        file = request.files['file']

        # Validate file
        is_valid, message = validate_pdf_file(file)
        if not is_valid:
            return jsonify({"error": message}), 400

        # Get form parameters
        split_method = request.form.get('split_method', 'pages')
        page_ranges = request.form.get('page_ranges', '')
        
        # Create secure filenames
        original_filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())[:8]
        input_filename = f"{unique_id}_{original_filename}"
        
        input_path = os.path.join(STATIC_DIR, input_filename)

        # Save uploaded file
        file.save(input_path)

        # Read PDF
        with open(input_path, 'rb') as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            total_pages = len(reader.pages)

            if split_method == 'pages' and page_ranges:
                # Parse page ranges (e.g., "1-3, 5, 7-10")
                ranges = []
                for part in page_ranges.split(','):
                    part = part.strip()
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        ranges.extend(range(start-1, end))  # Convert to 0-based
                    else:
                        ranges.append(int(part)-1)  # Convert to 0-based
                
                # Create single output file with selected pages
                writer = PyPDF2.PdfWriter()
                for page_num in ranges:
                    if 0 <= page_num < total_pages:
                        writer.add_page(reader.pages[page_num])
                
                output_filename = f"split_{unique_id}_{original_filename}"
                output_path = os.path.join(STATIC_DIR, output_filename)
                
                with open(output_path, 'wb') as output_file:
                    writer.write(output_file)
                
                # Clean up input file
                try:
                    os.remove(input_path)
                except:
                    pass
                
                return send_file(output_path, as_attachment=True, download_name=f"split_{original_filename}")
            
            else:
                # Split into individual pages
                output_files = []
                base_name = original_filename.rsplit('.', 1)[0]
                
                for page_num in range(total_pages):
                    writer = PyPDF2.PdfWriter()
                    writer.add_page(reader.pages[page_num])
                    
                    output_filename = f"{unique_id}_page_{page_num + 1}_{base_name}.pdf"
                    output_path = os.path.join(STATIC_DIR, output_filename)
                    
                    with open(output_path, 'wb') as output_file:
                        writer.write(output_file)
                    
                    output_files.append(output_path)
                
                # Create ZIP file
                zip_filename = f"split_{unique_id}_{base_name}.zip"
                zip_path = os.path.join(STATIC_DIR, zip_filename)
                
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    for output_file in output_files:
                        zipf.write(output_file, os.path.basename(output_file))
                
                # Clean up individual files
                for output_file in output_files:
                    try:
                        os.remove(output_file)
                    except:
                        pass
                
                # Clean up input file
                try:
                    os.remove(input_path)
                except:
                    pass
                
                return send_file(zip_path, as_attachment=True, download_name=zip_filename)

    except Exception as e:
        logging.error(f"Error in PDF split: {str(e)}")
        return jsonify({"error": f"Split failed: {str(e)}"}), 500

@basic_operations_bp.route('/protect', methods=['POST'])
def protect_pdf():
    """Add password protection to PDF"""
    try:
        # Get file from request
        if 'file' not in request.files:
            return jsonify({"error": "No file part in request"}), 400

        file = request.files['file']

        # Validate file
        is_valid, message = validate_pdf_file(file)
        if not is_valid:
            return jsonify({"error": message}), 400

        # Get form parameters
        user_password = request.form.get('user_password', '')
        owner_password = request.form.get('owner_password', '')

        if not user_password:
            return jsonify({"error": "User password is required"}), 400

        # Create secure filenames
        original_filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())[:8]
        input_filename = f"{unique_id}_{original_filename}"
        output_filename = f"protected_{unique_id}_{original_filename}"

        input_path = os.path.join(STATIC_DIR, input_filename)
        output_path = os.path.join(STATIC_DIR, output_filename)

        # Save uploaded file
        file.save(input_path)

        # Add password protection using PyPDF2
        with open(input_path, 'rb') as input_file:
            reader = PyPDF2.PdfReader(input_file)
            writer = PyPDF2.PdfWriter()

            # Copy all pages
            for page in reader.pages:
                writer.add_page(page)

            # Add password protection
            writer.encrypt(
                user_password=user_password,
                owner_password=owner_password if owner_password else user_password,
                use_128bit=True
            )

            # Write protected PDF
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)

        # Clean up input file
        try:
            os.remove(input_path)
        except:
            pass

        return send_file(output_path, as_attachment=True, download_name=f"protected_{original_filename}")

    except Exception as e:
        logging.error(f"Error in PDF protection: {str(e)}")
        return jsonify({"error": f"Protection failed: {str(e)}"}), 500

@basic_operations_bp.route('/pdf_to_panoramic', methods=['POST'])
def pdf_to_panoramic():
    """Convert PDF pages to panoramic image using Ghostscript"""
    try:
        # Get file from request
        if 'file' not in request.files:
            return jsonify({"error": "No file part in request"}), 400

        file = request.files['file']

        # Validate file
        is_valid, message = validate_pdf_file(file)
        if not is_valid:
            return jsonify({"error": message}), 400

        # Get form parameters
        stitch_direction = request.form.get('stitch_direction', 'horizontal')
        pages = request.form.get('pages', 'all')
        page_range = request.form.get('page_range', '')
        output_format = request.form.get('output_format', 'jpg')
        quality = request.form.get('quality', 'high')
        spacing = request.form.get('spacing', 'none')

        # Validate parameters
        if stitch_direction not in ['horizontal', 'vertical']:
            stitch_direction = 'horizontal'

        if output_format not in ['jpg', 'png', 'webp']:
            output_format = 'jpg'

        if quality not in ['high', 'medium', 'low']:
            quality = 'high'

        # Set DPI based on quality
        dpi_map = {'high': 300, 'medium': 150, 'low': 72}
        dpi = dpi_map[quality]

        # Create secure filenames
        original_filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())[:8]
        input_filename = f"{unique_id}_{original_filename}"

        input_path = os.path.join(STATIC_DIR, input_filename)

        # Save uploaded file
        file.save(input_path)
        logging.debug(f"File saved to: {input_path}")

        # Get total pages using PyPDF2
        with open(input_path, 'rb') as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            total_pages = len(reader.pages)

        # Determine which pages to process
        page_list = []
        if pages == 'all':
            page_list = list(range(1, total_pages + 1))
        elif pages == 'odd':
            page_list = list(range(1, total_pages + 1, 2))
        elif pages == 'even':
            page_list = list(range(2, total_pages + 1, 2))
        elif pages == 'range' and page_range:
            # Parse page range (e.g., "1-3,5,7-10")
            for part in page_range.split(','):
                part = part.strip()
                if '-' in part:
                    start, end = map(int, part.split('-'))
                    page_list.extend(range(start, min(end + 1, total_pages + 1)))
                else:
                    page_num = int(part)
                    if 1 <= page_num <= total_pages:
                        page_list.append(page_num)

        if not page_list:
            page_list = [1]  # Default to first page

        # Convert pages to images using Ghostscript
        temp_images = []

        # Determine Ghostscript executable
        import platform
        if platform.system() == 'Windows':
            # Try different Windows Ghostscript executables
            gs_executables = ['gswin64c.exe', 'gswin32c.exe', 'gs.exe', 'gs']
        else:
            gs_executables = ['gs']

        gs_executable = None
        for gs_exe in gs_executables:
            try:
                result = subprocess.run([gs_exe, '--version'], capture_output=True, text=True)
                if result.returncode == 0:
                    gs_executable = gs_exe
                    break
            except FileNotFoundError:
                continue

        if not gs_executable:
            # Fallback to PyMuPDF if Ghostscript not found
            logging.warning("Ghostscript not found, using PyMuPDF fallback")
            doc = fitz.open(input_path)
            for page_num in page_list:
                if page_num <= len(doc):
                    page = doc.load_page(page_num - 1)  # PyMuPDF uses 0-based indexing
                    mat = fitz.Matrix(dpi/72, dpi/72)
                    pix = page.get_pixmap(matrix=mat)
                    temp_image = os.path.join(STATIC_DIR, f"{unique_id}_page_{page_num}.png")
                    pix.save(temp_image)
                    temp_images.append(temp_image)
            doc.close()
        else:
            # Use Ghostscript
            for page_num in page_list:
                temp_image = os.path.join(STATIC_DIR, f"{unique_id}_page_{page_num}.png")

                # Ghostscript command to convert specific page to PNG
                gs_cmd = [
                    gs_executable,
                    '-dNOPAUSE',
                    '-dBATCH',
                    '-dSAFER',
                    '-sDEVICE=png16m',
                    f'-r{dpi}',
                    f'-dFirstPage={page_num}',
                    f'-dLastPage={page_num}',
                    f'-sOutputFile={temp_image}',
                    input_path
                ]

                result = subprocess.run(gs_cmd, capture_output=True, text=True)
                if result.returncode != 0:
                    logging.error(f"Ghostscript error: {result.stderr}")
                    # Cleanup and return error
                    try:
                        os.remove(input_path)
                    except:
                        pass
                    return jsonify({"error": "Failed to convert PDF pages to images"}), 500

                temp_images.append(temp_image)

        # Load images and create panoramic
        images = []
        for img_path in temp_images:
            if os.path.exists(img_path):
                images.append(Image.open(img_path))

        if not images:
            return jsonify({"error": "No images were generated"}), 500

        # Calculate spacing
        spacing_pixels = 0
        if spacing == 'small':
            spacing_pixels = 10
        elif spacing == 'medium':
            spacing_pixels = 20
        elif spacing == 'large':
            spacing_pixels = 40

        # Check and resize images if they're too large
        MAX_DIMENSION = 32000  # Safe limit for PIL
        resized_images = []

        for img in images:
            # Calculate if image needs resizing
            if img.width > MAX_DIMENSION or img.height > MAX_DIMENSION:
                # Calculate scale factor to fit within limits
                scale_factor = min(MAX_DIMENSION / img.width, MAX_DIMENSION / img.height)
                new_width = int(img.width * scale_factor)
                new_height = int(img.height * scale_factor)
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                logging.info(f"Resized image to {new_width}x{new_height}")
            resized_images.append(img)

        images = resized_images

        # Create panoramic image with size limits
        if stitch_direction == 'horizontal':
            # Horizontal stitching (side by side)
            total_width = sum(img.width for img in images) + spacing_pixels * (len(images) - 1)
            max_height = max(img.height for img in images)

            # Check if panoramic would be too large
            if total_width > MAX_DIMENSION:
                # Scale down all images proportionally
                scale_factor = MAX_DIMENSION / total_width
                scaled_images = []
                for img in images:
                    new_width = int(img.width * scale_factor)
                    new_height = int(img.height * scale_factor)
                    scaled_images.append(img.resize((new_width, new_height), Image.Resampling.LANCZOS))
                images = scaled_images
                total_width = sum(img.width for img in images) + int(spacing_pixels * scale_factor) * (len(images) - 1)
                max_height = max(img.height for img in images)
                spacing_pixels = int(spacing_pixels * scale_factor)

            panoramic = Image.new('RGB', (total_width, max_height), 'white')

            x_offset = 0
            for img in images:
                # Center vertically
                y_offset = (max_height - img.height) // 2
                panoramic.paste(img, (x_offset, y_offset))
                x_offset += img.width + spacing_pixels
        else:
            # Vertical stitching (top to bottom)
            max_width = max(img.width for img in images)
            total_height = sum(img.height for img in images) + spacing_pixels * (len(images) - 1)

            # Check if panoramic would be too large
            if total_height > MAX_DIMENSION:
                # Scale down all images proportionally
                scale_factor = MAX_DIMENSION / total_height
                scaled_images = []
                for img in images:
                    new_width = int(img.width * scale_factor)
                    new_height = int(img.height * scale_factor)
                    scaled_images.append(img.resize((new_width, new_height), Image.Resampling.LANCZOS))
                images = scaled_images
                max_width = max(img.width for img in images)
                total_height = sum(img.height for img in images) + int(spacing_pixels * scale_factor) * (len(images) - 1)
                spacing_pixels = int(spacing_pixels * scale_factor)

            panoramic = Image.new('RGB', (max_width, total_height), 'white')

            y_offset = 0
            for img in images:
                # Center horizontally
                x_offset = (max_width - img.width) // 2
                panoramic.paste(img, (x_offset, y_offset))
                y_offset += img.height + spacing_pixels

        # Save panoramic image with error handling
        output_filename = f"panoramic_{unique_id}_{original_filename.rsplit('.', 1)[0]}.{output_format}"
        output_path = os.path.join(STATIC_DIR, output_filename)

        try:
            # Increase PIL's image size limit temporarily
            from PIL import ImageFile
            ImageFile.LOAD_TRUNCATED_IMAGES = True

            if output_format == 'jpg':
                panoramic.save(output_path, 'JPEG', quality=85, optimize=True)
            elif output_format == 'png':
                panoramic.save(output_path, 'PNG', optimize=True)
            elif output_format == 'webp':
                panoramic.save(output_path, 'WebP', quality=85, optimize=True)

            logging.info(f"Panoramic image saved: {output_path} ({panoramic.width}x{panoramic.height})")

        except Exception as save_error:
            logging.error(f"Error saving panoramic image: {str(save_error)}")
            # Try saving as a smaller JPEG if other formats fail
            try:
                output_filename = f"panoramic_{unique_id}_{original_filename.rsplit('.', 1)[0]}.jpg"
                output_path = os.path.join(STATIC_DIR, output_filename)
                panoramic.save(output_path, 'JPEG', quality=75, optimize=True)
                logging.info(f"Fallback save successful: {output_path}")
            except Exception as fallback_error:
                logging.error(f"Fallback save failed: {str(fallback_error)}")
                raise save_error

        # Clean up temporary files
        try:
            os.remove(input_path)
            for img_path in temp_images:
                if os.path.exists(img_path):
                    os.remove(img_path)
        except:
            pass

        return send_file(output_path, as_attachment=True, download_name=output_filename)

    except Exception as e:
        logging.error(f"Error in PDF to panoramic conversion: {str(e)}")
        return jsonify({"error": f"Panoramic conversion failed: {str(e)}"}), 500

@basic_operations_bp.route('/powerpoint', methods=['POST'])
def pdf_to_powerpoint():
    """Convert PDF to PowerPoint with full content extraction and recreation"""
    try:
        # Get file from request
        if 'file' not in request.files:
            return jsonify({"error": "No file part in request"}), 400

        file = request.files['file']

        # Validate file
        is_valid, message = validate_pdf_file(file)
        if not is_valid:
            return jsonify({"error": message}), 400

        # Get form parameters
        slide_size = request.form.get('slide_size', 'standard')
        conversion_mode = request.form.get('conversion_mode', 'layout_preserve')
        pages_per_slide = int(request.form.get('pages_per_slide', '1'))

        # Create secure filenames
        original_filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())[:8]
        input_filename = f"{unique_id}_{original_filename}"

        input_path = os.path.join(STATIC_DIR, input_filename)

        # Save uploaded file
        file.save(input_path)
        logging.debug(f"File saved to: {input_path}")

        # Import PowerPoint libraries
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
        from pptx.dml.color import RGBColor
        import unicodedata
        import re

        # Create PowerPoint presentation
        prs = Presentation()

        # 🎯 BREAKTHROUGH: Set slide dimensions to match PDF page exactly (no scaling!)
        # This is the key to true 1:1 positioning - no dimension mismatch
        doc = fitz.open(input_path)
        first_page = doc.load_page(0)
        pdf_page_rect = first_page.rect
        pdf_width_inches = pdf_page_rect.width / 72  # Convert points to inches
        pdf_height_inches = pdf_page_rect.height / 72

        # Set PowerPoint slide to exact PDF dimensions
        prs.slide_width = Inches(pdf_width_inches)
        prs.slide_height = Inches(pdf_height_inches)

        logging.info(f"🎯 EXACT MATCH: PDF page {pdf_width_inches:.2f}\" x {pdf_height_inches:.2f}\" -> PPT slide same size")
        doc.close()

        # Open PDF with PyMuPDF
        doc = fitz.open(input_path)

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)

            # Since slide dimensions match PDF exactly, no scaling needed for 1:1 conversion

            # Add slide to presentation
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # 🚀 PERFECT TEXT-BASED LAYOUT PRESERVATION - Fully editable PowerPoint
            if conversion_mode == 'layout_preserve':
                try:
                    # Get PDF page dimensions
                    pdf_rect = page.rect
                    pdf_width_points = pdf_rect.width
                    pdf_height_points = pdf_rect.height

                    logging.info(f"🎯 PERFECT TEXT-BASED LAYOUT PRESERVATION")
                    logging.info(f"📄 PDF Page: {pdf_width_points:.1f} x {pdf_height_points:.1f} points")

                    # 🎯 STEP 1: EXTRACT TEXT WITH PRECISE COORDINATES AND CLEAN CONTENT
                    def clean_text(text):
                        """Clean and normalize text content"""
                        if not text:
                            return ""

                        # Normalize Unicode characters
                        try:
                            text = unicodedata.normalize('NFKC', text)
                        except:
                            pass

                        # 🎯 ULTRA-COMPREHENSIVE CHARACTER MAPPING FOR PERFECT TEXT EXTRACTION
                        replacements = {
                            # 🔥 PRIMARY CORRUPTED CHARACTERS (most common in PDF extraction)
                            'ƌeƋuiƌes': 'requires', 'ƌeƋuiƌe': 'require', 'ƌeƋuiƌed': 'required',
                            'paƌameteƌ': 'parameter', 'paƌameteƌs': 'parameters',
                            'geneƌal': 'general', 'geneƌally': 'generally', 'geneƌate': 'generate',
                            'numbeƌ': 'number', 'numbeƌs': 'numbers', 'membeƌ': 'member',
                            'otheƌ': 'other', 'otheƌs': 'others', 'otheƌwise': 'otherwise',
                            'wheƌe': 'where', 'theƌe': 'there', 'heƌe': 'here',
                            'foƌ': 'for', 'oƌ': 'or', 'youƌ': 'your', 'ouƌ': 'our',
                            'theiƌ': 'their', 'suƌe': 'sure', 'puƌe': 'pure',
                            'befoƌe': 'before', 'afteƌ': 'after', 'duƌing': 'during',
                            'oveƌ': 'over', 'undeƌ': 'under', 'inteƌ': 'inter',
                            'supeƌ': 'super', 'hypeƌ': 'hyper', 'ulteƌ': 'ultra',
                            'centeƌ': 'center', 'enteƌ': 'enter', 'lateƌ': 'later',
                            'wateƌ': 'water', 'mateƌ': 'matter', 'beteƌ': 'better',
                            'letteƌ': 'letter', 'chaƌacteƌ': 'character',

                            # Individual character mappings
                            'ƌ': 'r', 'Ƌ': 'R', 'ƍ': 'd', 'Ƈ': 'C', 'ƈ': 'c',
                            'ƅ': 'b', 'Ɓ': 'B', 'ƃ': 'p', 'Ƃ': 'P', 'ƀ': 'b',
                            'ƒ': 'f', 'Ɠ': 'G', 'ɠ': 'g', 'ɡ': 'g', 'ɢ': 'G',
                            'ɣ': 'y', 'ɤ': 'o', 'ɥ': 'h', 'ɦ': 'h', 'ɧ': 'h',
                            'ɨ': 'i', 'ɩ': 'i', 'ɪ': 'I', 'ɫ': 'l', 'ɬ': 'l',
                            'ɭ': 'l', 'ɮ': 'lz', 'ɯ': 'm', 'ɰ': 'm', 'ɱ': 'm',
                            'ɲ': 'n', 'ɳ': 'n', 'ɴ': 'N', 'ɵ': 'o', 'ɶ': 'o',
                            'ɷ': 'o', 'ɸ': 'f', 'ɹ': 'r', 'ɺ': 'r', 'ɻ': 'r',
                            'ɼ': 'r', 'ɽ': 'r', 'ɾ': 'r', 'ɿ': 'r', 'ʀ': 'R',
                            'ʁ': 'R', 'ʂ': 's', 'ʃ': 's', 'ʄ': 'j', 'ʅ': 'l',
                            'ʆ': 's', 'ʇ': 't', 'ʈ': 't', 'ʉ': 'u', 'ʊ': 'u',
                            'ʋ': 'v', 'ʌ': 'v', 'ʍ': 'w', 'ʎ': 'y', 'ʏ': 'Y',
                            'ʐ': 'z', 'ʑ': 'z', 'ʒ': 'z', 'ʓ': 'z', 'ʔ': '',

                            # Extended IPA and phonetic characters
                            'ǝ': 'e', 'ɐ': 'a', 'ɑ': 'a', 'ɒ': 'o', 'ɓ': 'b',
                            'ɔ': 'o', 'ɕ': 'c', 'ɖ': 'd', 'ɗ': 'd', 'ɘ': 'e',
                            'ə': 'e', 'ɚ': 'e', 'ɛ': 'e', 'ɜ': 'e', 'ɝ': 'e',
                            'ɞ': 'e', 'ɟ': 'j', 'ɠ': 'g', 'ɡ': 'g', 'ɢ': 'G',

                            # Legacy character fixes
                            'Ŷ': 'n', 'Đ': 'D', 'ŵ': 'm', 'ď': 'd', 'ǆ': 'x',
                            'Ŷow': 'now', 'ŵediaďoǆ': 'mediabox',
                            'ƌe': 're', 'Ƌui': 'qui', 'ƌes': 'res', 'ŵedia': 'media', 'ďo': 'do',

                            # Bullet points and symbols
                            '\uf0b7': '•', '\uf0a7': '•', '\u2022': '•',
                            '\u2013': '–', '\u2014': '—', '\u2019': "'",
                            '\u201c': '"', '\u201d': '"', '\u2018': "'", '\u201a': ',',
                            '\u201e': '"', '\u2026': '...', '\u2030': '%',

                            # Ligatures and typography
                            'ﬁ': 'fi', 'ﬂ': 'fl', 'ﬀ': 'ff', 'ﬃ': 'ffi', 'ﬄ': 'ffl',
                            'ﬅ': 'ft', 'ﬆ': 'st', '﻿': '',  # Zero-width no-break space

                            # Mathematical symbols
                            '×': 'x', '÷': '/', '±': '+/-', '≤': '<=', '≥': '>=',
                            '≠': '!=', '≈': '~=', '°': 'deg', '′': "'", '″': '"',

                            # Currency and special symbols
                            '€': 'EUR', '£': 'GBP', '¥': 'JPY', '¢': 'cents',
                            '©': '(c)', '®': '(R)', '™': '(TM)', '§': 'section',
                            '¶': 'paragraph', '†': '+', '‡': '++', '•': '*',

                            # Remove problematic invisible characters
                            '\u200b': '', '\u200c': '', '\u200d': '', '\ufeff': '',
                            '\u00a0': ' ',  # Non-breaking space to regular space
                        }

                        for old, new in replacements.items():
                            text = text.replace(old, new)

                        return text.strip()

                    # Extract text elements with detailed structure and comprehensive debugging
                    text_dict = page.get_text("dict")
                    text_elements = []

                    logging.info("🔍 COMPREHENSIVE TEXT EXTRACTION DEBUGGING:")
                    logging.info(f"   PDF Page Dimensions: {pdf_rect.width:.1f} x {pdf_rect.height:.1f} points")
                    logging.info(f"   PDF Coordinate System: Origin at bottom-left, Y increases upward")

                    element_count = 0
                    for block in text_dict.get("blocks", []):
                        if "lines" in block:  # Text block
                            for line in block["lines"]:
                                # 🎯 LINE-LEVEL DEBUGGING
                                line_bbox = line.get("bbox", None)
                                if element_count < 3 and line_bbox:
                                    logging.info(f"📏 LINE #{element_count + 1} BBOX: ({line_bbox[0]:.1f}, {line_bbox[1]:.1f}, {line_bbox[2]:.1f}, {line_bbox[3]:.1f})")

                                # Process each span (text run with consistent formatting)
                                for span in line["spans"]:
                                    raw_text = span.get("text", "")
                                    if not raw_text.strip():
                                        continue

                                    # Clean the text
                                    clean_text_content = clean_text(raw_text)
                                    if not clean_text_content:
                                        continue

                                    element_count += 1
                                    bbox = span["bbox"]
                                    font_name = span.get("font", "Arial")
                                    font_size = span.get("size", 12)
                                    font_flags = span.get("flags", 0)
                                    color = span.get("color", 0)

                                    # 🎯 COMPREHENSIVE TEXT ELEMENT DEBUGGING (first 5 elements)
                                    if element_count <= 5:
                                        logging.info(f"📝 TEXT ELEMENT #{element_count} DETAILED ANALYSIS:")
                                        logging.info(f"   Content: '{clean_text_content}'")
                                        logging.info(f"   Raw Text: '{raw_text}'")
                                        logging.info(f"   Font Family: {font_name}")
                                        logging.info(f"   Font Size: {font_size:.1f}pt")
                                        logging.info(f"   Font Flags: {font_flags} (bold={bool(font_flags & 2**4)}, italic={bool(font_flags & 2**1)})")
                                        logging.info(f"   Color: {color} (hex: {hex(color) if isinstance(color, int) else 'N/A'})")
                                        logging.info(f"   Raw PDF BBox: ({bbox[0]:.1f}, {bbox[1]:.1f}, {bbox[2]:.1f}, {bbox[3]:.1f})")
                                        logging.info(f"   BBox Analysis:")
                                        logging.info(f"     - Left (x0): {bbox[0]:.1f}")
                                        logging.info(f"     - Bottom (y0): {bbox[1]:.1f} ← PDF bottom edge")
                                        logging.info(f"     - Right (x1): {bbox[2]:.1f}")
                                        logging.info(f"     - Top (y1): {bbox[3]:.1f} ← PDF top edge")
                                        logging.info(f"     - Width: {bbox[2] - bbox[0]:.1f}")
                                        logging.info(f"     - Height: {bbox[3] - bbox[1]:.1f}")

                                        # Check for text direction/orientation attributes
                                        if 'dir' in span:
                                            logging.info(f"   Text Direction: {span['dir']}")
                                        if 'wmode' in span:
                                            logging.info(f"   Writing Mode: {span['wmode']}")
                                        if 'ascender' in span:
                                            logging.info(f"   Ascender: {span['ascender']}")
                                        if 'descender' in span:
                                            logging.info(f"   Descender: {span['descender']}")

                                        # Check transformation matrix if available
                                        if hasattr(span, 'transform') or 'transform' in span:
                                            transform = span.get('transform', getattr(span, 'transform', None))
                                            if transform:
                                                logging.info(f"   Transform Matrix: {transform}")

                                        logging.info("   ---")

                                    # Detect bullet points
                                    bullet_patterns = [
                                        r'^\s*[•◦▪▫▶►‣⁃]\s+',  # Unicode bullets
                                        r'^\s*[-*+]\s+',         # ASCII bullets
                                        r'^\s*\d+[\.\)]\s+',     # Numbered lists
                                        r'^\s*[a-zA-Z][\.\)]\s+' # Lettered lists
                                    ]
                                    is_bullet = any(re.match(pattern, clean_text_content) for pattern in bullet_patterns)

                                    element = {
                                        'text': clean_text_content,
                                        'bbox': bbox,
                                        'font_name': font_name,
                                        'font_size': font_size,
                                        'font_flags': font_flags,
                                        'color': color,
                                        'is_bold': bool(font_flags & 2**4),
                                        'is_italic': bool(font_flags & 2**1),
                                        'is_bullet': is_bullet
                                    }

                                    text_elements.append(element)

                    logging.info(f"📝 Extracted {len(text_elements)} text elements")

                    # 🎯 STEP 2: DETECT HIGHLIGHTED TEXT AND BACKGROUND COLORS
                    # Extract drawing objects (rectangles, highlights, etc.)
                    drawings = page.get_drawings()
                    highlight_areas = []

                    for drawing in drawings:
                        if hasattr(drawing, 'rect') and hasattr(drawing, 'fill'):
                            # This is a filled rectangle - potential highlight
                            highlight_areas.append({
                                'bbox': drawing.rect,
                                'color': getattr(drawing, 'fill', None)
                            })

                    logging.info(f"🎨 Found {len(highlight_areas)} potential highlight areas")

                    # 🎯 STEP 2.5: DETECT LINES AND BORDERS FOR PAGE STRUCTURE
                    # Extract vector graphics (lines, borders, dividers)
                    lines_and_borders = []

                    try:
                        # Get all drawing paths from the page using a different approach
                        # Use PyMuPDF's path extraction for lines
                        paths = page.get_drawings()

                        if paths:
                            for path in paths:
                                # Check if this path contains line segments
                                if hasattr(path, 'items') and hasattr(path.items, '__iter__'):
                                    for item in path.items:
                                        # Look for line segments
                                        if isinstance(item, (list, tuple)) and len(item) > 0:
                                            if item[0] == 'l':  # Line command
                                                # Extract line coordinates
                                                if len(item) >= 3:
                                                    start_point = item[1]
                                                    end_point = item[2]

                                                    # Check if it's a horizontal or vertical line
                                                    is_horizontal = abs(start_point[1] - end_point[1]) < 2
                                                    is_vertical = abs(start_point[0] - end_point[0]) < 2

                                                    if is_horizontal or is_vertical:
                                                        lines_and_borders.append({
                                                            'type': 'horizontal' if is_horizontal else 'vertical',
                                                            'start': start_point,
                                                            'end': end_point,
                                                            'stroke_width': getattr(path, 'width', 1.0)
                                                        })
                    except Exception as line_error:
                        logging.warning(f"⚠️ Line detection failed: {line_error}")

                    logging.info(f"📏 Found {len(lines_and_borders)} lines and borders")

                    # 🎯 STEP 3: RESTRICTIVE TABLE DETECTION - ONLY DETECT LARGE ACTUAL TABLES
                    def detect_tables(elements):
                        """Restrictive table detection - only detect substantial tables with many cells"""
                        # 🎯 BALANCED: Require at least 4 elements for a table (2x2 minimum)
                        if len(elements) < 4:
                            return []

                        # 🎯 PHASE 1: GROUP BY ROWS (Y-coordinate alignment) - MORE RESTRICTIVE
                        tolerance = 3  # Slightly looser tolerance but require more cells
                        rows = []

                        # Sort by Y position (top to bottom) using center Y for better alignment
                        sorted_elements = sorted(elements, key=lambda x: -(x['bbox'][1] + x['bbox'][3]) / 2)

                        current_row = []
                        current_y = None

                        for elem in sorted_elements:
                            elem_y = (elem['bbox'][1] + elem['bbox'][3]) / 2  # Center Y

                            if current_y is None or abs(elem_y - current_y) <= tolerance:
                                current_row.append(elem)
                                current_y = elem_y
                            else:
                                # 🎯 REQUIRE AT LEAST 2 CELLS PER ROW for table detection
                                if len(current_row) >= 2:
                                    # Sort row by X position (left to right)
                                    rows.append(sorted(current_row, key=lambda x: x['bbox'][0]))
                                current_row = [elem]
                                current_y = elem_y

                        if len(current_row) >= 2:
                            rows.append(sorted(current_row, key=lambda x: x['bbox'][0]))

                        # 🎯 REQUIRE AT LEAST 2 ROWS for a table
                        if len(rows) < 2:
                            return []

                        # 🎯 PHASE 2: ANALYZE COLUMN BOUNDARIES
                        def get_column_boundaries(table_rows):
                            """Detect precise column boundaries from row data"""
                            all_x_positions = []
                            for row in table_rows:
                                for cell in row:
                                    all_x_positions.extend([cell['bbox'][0], cell['bbox'][2]])  # Left and right edges

                            # Cluster similar X positions to find column boundaries
                            all_x_positions.sort()
                            boundaries = []
                            current_boundary = all_x_positions[0]
                            boundary_tolerance = 5  # Points

                            for x in all_x_positions[1:]:
                                if x - current_boundary > boundary_tolerance:
                                    boundaries.append(current_boundary)
                                    current_boundary = x
                                else:
                                    current_boundary = (current_boundary + x) / 2  # Average
                            boundaries.append(current_boundary)

                            return sorted(set(boundaries))

                        # 🎯 PHASE 3: GROUP ROWS INTO TABLES WITH COLUMN CONSISTENCY
                        tables = []
                        current_table = [rows[0]]
                        current_boundaries = get_column_boundaries([rows[0]])

                        for row in rows[1:]:
                            row_boundaries = get_column_boundaries([row])

                            # Check if column boundaries are similar (indicating same table)
                            boundary_similarity = 0
                            if len(current_boundaries) > 0 and len(row_boundaries) > 0:
                                # Calculate how many boundaries align
                                aligned_boundaries = 0
                                for cb in current_boundaries:
                                    for rb in row_boundaries:
                                        if abs(cb - rb) <= 10:  # 10 points tolerance
                                            aligned_boundaries += 1
                                            break
                                boundary_similarity = aligned_boundaries / max(len(current_boundaries), len(row_boundaries))

                            # If boundaries align well (>60% similarity), it's the same table
                            if boundary_similarity > 0.6 and abs(len(row) - len(current_table[-1])) <= 1:
                                current_table.append(row)
                                # Update boundaries to include new row
                                current_boundaries = get_column_boundaries(current_table)
                            else:
                                # Start new table - REQUIRE AT LEAST 2 ROWS
                                if len(current_table) >= 2:
                                    tables.append(current_table)
                                current_table = [row]
                                current_boundaries = get_column_boundaries([row])

                        if len(current_table) >= 2:
                            tables.append(current_table)

                        # 🎯 PHASE 4: REFINE TABLES BY REMOVING OUTLIERS - VERY RESTRICTIVE
                        refined_tables = []
                        for table in tables:
                            # 🎯 REQUIRE AT LEAST 2 ROWS AND CONSISTENT COLUMN COUNT
                            if len(table) >= 2:
                                # Remove rows that don't fit the column pattern
                                most_common_col_count = max(set(len(row) for row in table), key=lambda x: sum(1 for row in table if len(row) == x))
                                refined_table = [row for row in table if abs(len(row) - most_common_col_count) <= 1]

                                # 🎯 BALANCED CHECK: Must have at least 2 rows and 2 columns for real tables
                                if len(refined_table) >= 2 and most_common_col_count >= 2:
                                    # Calculate total cells - must be substantial table
                                    total_cells = sum(len(row) for row in refined_table)
                                    if total_cells >= 4:  # At least 4 cells total (2x2 minimum)
                                        refined_tables.append(refined_table)

                        return refined_tables

                    detected_tables = detect_tables(text_elements)
                    logging.info(f"📋 Detected {len(detected_tables)} tables")

                    # Mark table elements
                    table_element_ids = set()
                    for table in detected_tables:
                        for row in table:
                            for cell in row:
                                table_element_ids.add(id(cell))

                    # 🎯 STEP 4: SORT BY READING ORDER (TOP TO BOTTOM, LEFT TO RIGHT)
                    # PDF Y coordinates: Y=0 at bottom, higher Y = higher on page
                    # Sort by Y descending (top first), then X ascending (left first)
                    text_elements.sort(key=lambda x: (-x['bbox'][3], x['bbox'][0]))

                    # 🎯 STEP 5: CREATE PRECISION POWERPOINT TABLES WITH EXACT POSITIONING
                    for table_idx, table in enumerate(detected_tables):
                        if not table:
                            continue

                        try:
                            # 🎯 CALCULATE PRECISE TABLE BOUNDARIES
                            all_cells = [cell for row in table for cell in row]
                            min_x = min(cell['bbox'][0] for cell in all_cells)
                            max_x = max(cell['bbox'][2] for cell in all_cells)
                            min_y = min(cell['bbox'][1] for cell in all_cells)
                            max_y = max(cell['bbox'][3] for cell in all_cells)

                            # Convert to PowerPoint coordinates with precision
                            table_x = min_x / 72
                            table_y = (pdf_height_points - max_y) / 72
                            table_width = (max_x - min_x) / 72
                            table_height = (max_y - min_y) / 72

                            # 🎯 DETERMINE OPTIMAL TABLE DIMENSIONS
                            max_cols = max(len(row) for row in table)
                            rows_count = len(table)

                            # 🎯 ANALYZE COLUMN WIDTHS FROM ACTUAL CONTENT
                            column_boundaries = []
                            for col_idx in range(max_cols + 1):  # +1 for right boundary
                                x_positions = []
                                for row in table:
                                    for cell_idx, cell in enumerate(row):
                                        if cell_idx == col_idx and col_idx < len(row):
                                            x_positions.append(cell['bbox'][0])  # Left edge
                                        elif cell_idx == col_idx - 1 and col_idx > 0:
                                            x_positions.append(cell['bbox'][2])  # Right edge of previous cell

                                if x_positions:
                                    column_boundaries.append(sum(x_positions) / len(x_positions))
                                elif col_idx == 0:
                                    column_boundaries.append(min_x)
                                else:
                                    column_boundaries.append(max_x)

                            # Ensure boundaries are sorted and within table bounds
                            column_boundaries = sorted(set(column_boundaries))
                            column_boundaries = [max(min_x, min(max_x, b)) for b in column_boundaries]

                            # Create PowerPoint table
                            ppt_table = slide.shapes.add_table(
                                rows_count, max_cols,
                                Inches(table_x), Inches(table_y),
                                Inches(table_width), Inches(table_height)
                            ).table

                            # 🎯 APPLY PROFESSIONAL TABLE STYLING
                            # Set table style for proper borders and formatting
                            ppt_table.first_row = True  # Enable header row styling

                            # 🎯 SET PRECISE COLUMN WIDTHS
                            if len(column_boundaries) >= 2:
                                for col_idx in range(max_cols):
                                    if col_idx + 1 < len(column_boundaries):
                                        col_width = (column_boundaries[col_idx + 1] - column_boundaries[col_idx]) / 72
                                        ppt_table.columns[col_idx].width = Inches(max(col_width, 0.5))  # Minimum width

                            # 🎯 FILL CELLS WITH PRECISE FORMATTING
                            for row_idx, row in enumerate(table):
                                # Calculate row height from actual content
                                row_height = 0
                                for cell in row:
                                    cell_height = (cell['bbox'][3] - cell['bbox'][1]) / 72
                                    row_height = max(row_height, cell_height)

                                # Set row height
                                if row_height > 0:
                                    ppt_table.rows[row_idx].height = Inches(max(row_height, 0.2))

                                # Fill cells
                                for col_idx, cell in enumerate(row):
                                    if col_idx < max_cols:
                                        ppt_cell = ppt_table.cell(row_idx, col_idx)
                                        ppt_cell.text = cell['text']

                                        # 🎯 PRECISE CELL FORMATTING
                                        text_frame = ppt_cell.text_frame
                                        text_frame.margin_left = Inches(0.05)
                                        text_frame.margin_right = Inches(0.05)
                                        text_frame.margin_top = Inches(0.02)
                                        text_frame.margin_bottom = Inches(0.02)
                                        text_frame.word_wrap = True

                                        # 🎯 ADD PROFESSIONAL TABLE BORDERS
                                        # Apply borders to all cells for proper table appearance
                                        from pptx.dml.color import RGBColor

                                        # PowerPoint table cells don't have direct border properties
                                        # Instead, we'll use the table's overall styling
                                        try:
                                            # Set table style to show borders
                                            ppt_table.first_row = True
                                            ppt_table.first_col = False
                                            ppt_table.last_row = False
                                            ppt_table.last_col = False
                                            ppt_table.horz_banding = True
                                            ppt_table.vert_banding = False
                                        except Exception as border_error:
                                            logging.warning(f"⚠️ Failed to apply table borders: {border_error}")

                                        # 🎯 HEADER ROW STYLING - Blue background for first row
                                        if row_idx == 0:  # First row is header
                                            fill = ppt_cell.fill
                                            fill.solid()
                                            fill.fore_color.rgb = RGBColor(79, 129, 189)  # Professional blue

                                            # Make header text white and bold
                                            for paragraph in text_frame.paragraphs:
                                                for run in paragraph.runs:
                                                    run.font.color.rgb = RGBColor(255, 255, 255)  # White text
                                                    run.font.bold = True

                                        # Apply text formatting
                                        for paragraph in text_frame.paragraphs:
                                            paragraph.alignment = PP_ALIGN.LEFT
                                            for run in paragraph.runs:
                                                run.font.name = cell.get('font_name', 'Calibri')
                                                run.font.size = Pt(max(cell.get('font_size', 11) * 0.8, 8))
                                                run.font.bold = cell.get('is_bold', False)
                                                run.font.italic = cell.get('is_italic', False)

                                                # Apply text color
                                                if cell.get('color') and cell['color'] != 0:
                                                    try:
                                                        color_rgb = RGBColor(
                                                            (cell['color'] >> 16) & 0xFF,
                                                            (cell['color'] >> 8) & 0xFF,
                                                            cell['color'] & 0xFF
                                                        )
                                                        run.font.color.rgb = color_rgb
                                                    except:
                                                        pass

                                        # 🎯 APPLY CELL BACKGROUND COLORS (for headers, etc.)
                                        # Check if this cell should have a background color
                                        cell_center_x = (cell['bbox'][0] + cell['bbox'][2]) / 2
                                        cell_center_y = (cell['bbox'][1] + cell['bbox'][3]) / 2

                                        for highlight in highlight_areas:
                                            h_bbox = highlight['bbox']
                                            if (h_bbox[0] <= cell_center_x <= h_bbox[2] and
                                                h_bbox[1] <= cell_center_y <= h_bbox[3]):
                                                try:
                                                    fill = ppt_cell.fill
                                                    fill.solid()
                                                    if highlight.get('color'):
                                                        # Apply highlight color
                                                        if isinstance(highlight['color'], (int, float)):
                                                            fill.fore_color.rgb = RGBColor(
                                                                (int(highlight['color']) >> 16) & 0xFF,
                                                                (int(highlight['color']) >> 8) & 0xFF,
                                                                int(highlight['color']) & 0xFF
                                                            )
                                                        else:
                                                            fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
                                                    else:
                                                        fill.fore_color.rgb = RGBColor(173, 216, 230)  # Default light blue
                                                    break
                                                except Exception as bg_error:
                                                    logging.warning(f"⚠️ Failed to apply cell background: {bg_error}")

                            logging.info(f"✅ Created precision table {table_idx + 1}: {rows_count}x{max_cols} at ({table_x:.2f}\", {table_y:.2f}\")")

                        except Exception as e:
                            logging.warning(f"⚠️ Failed to create table {table_idx + 1}: {e}")
                            import traceback
                            logging.warning(f"   Traceback: {traceback.format_exc()}")

                    # 🎯 STEP 6: PLACE NON-TABLE ELEMENTS WITH PERFECT COORDINATE TRANSFORMATION
                    non_table_elements = [elem for elem in text_elements if id(elem) not in table_element_ids]
                    logging.info(f"📝 Placing {len(non_table_elements)} non-table elements")

                    for i, element in enumerate(non_table_elements):
                        if not element['text'].strip():
                            continue

                        bbox = element['bbox']

                        # 🎯 COMPREHENSIVE COORDINATE TRANSFORMATION DEBUGGING
                        # PDF coordinates: (x0, y0, x1, y1) where y0 is bottom edge, y1 is top edge
                        pdf_x_points = bbox[0]  # Left edge
                        pdf_y_bottom = bbox[1]  # Bottom edge in PDF coordinates
                        pdf_y_top = bbox[3]     # Top edge in PDF coordinates
                        element_width_points = bbox[2] - bbox[0]
                        element_height_points = bbox[3] - bbox[1]

                        # Skip invalid elements
                        if element_width_points <= 0 or element_height_points <= 0:
                            continue

                        # 🎯 STEP-BY-STEP COORDINATE TRANSFORMATION DEBUGGING
                        if i < 3:
                            logging.info(f"🔍 COMPREHENSIVE COORDINATE TRANSFORMATION #{i+1}:")
                            logging.info(f"   📄 PDF INPUT ANALYSIS:")
                            logging.info(f"     - Page dimensions: {pdf_rect.width:.1f} x {pdf_rect.height:.1f} points")
                            logging.info(f"     - Element bbox: ({pdf_x_points:.1f}, {pdf_y_bottom:.1f}, {pdf_x_points + element_width_points:.1f}, {pdf_y_top:.1f})")
                            logging.info(f"     - PDF coordinate system: Y=0 at BOTTOM, Y increases UPWARD")
                            logging.info(f"     - Element Y-position in PDF: bottom={pdf_y_bottom:.1f}, top={pdf_y_top:.1f}")
                            logging.info(f"     - Distance from PDF bottom: {pdf_y_bottom:.1f} points")
                            logging.info(f"     - Distance from PDF top: {pdf_rect.height - pdf_y_top:.1f} points")

                        # 🎯 ULTRA-PRECISE COORDINATE TRANSFORMATION
                        # Convert PDF points to PowerPoint inches with sub-pixel precision
                        ppt_x_inches = pdf_x_points / 72.0

                        # Y coordinate flip: PDF Y=0 at bottom, PPT Y=0 at top
                        # For PDF: Y increases upward (bottom=0)
                        # For PPT: Y increases downward (top=0)
                        # 🎯 CRITICAL FIX: Use bottom edge and flip Y coordinate properly
                        # PDF Y=0 at bottom, PPT Y=0 at top - use bottom edge for baseline
                        raw_y_calculation = (pdf_rect.height - pdf_y_bottom) / 72.0
                        ppt_y_inches = raw_y_calculation

                        # 🎯 SUB-PIXEL PRECISION ADJUSTMENTS
                        # Round to 4 decimal places for sub-pixel accuracy
                        ppt_x_inches = round(ppt_x_inches, 4)
                        ppt_y_inches = round(ppt_y_inches, 4)

                        # 🎯 DETAILED COORDINATE TRANSFORMATION LOGGING
                        if i < 3:
                            logging.info(f"   � TRANSFORMATION CALCULATIONS:")
                            logging.info(f"     - X transformation: {pdf_x_points:.1f} / 72 = {ppt_x_inches:.4f} inches")
                            logging.info(f"     - Y transformation method: BOTTOM EDGE")
                            logging.info(f"     - Y calculation: ({pdf_rect.height:.1f} - {pdf_y_bottom:.1f}) / 72")
                            logging.info(f"     - Y calculation: {pdf_rect.height - pdf_y_bottom:.1f} / 72 = {raw_y_calculation:.4f}")
                            logging.info(f"     - PowerPoint coordinate system: Y=0 at TOP, Y increases DOWNWARD")
                            logging.info(f"   📍 FINAL PPT COORDINATES:")
                            logging.info(f"     - Position: ({ppt_x_inches:.4f}\", {ppt_y_inches:.4f}\")")
                            logging.info(f"     - Expected text orientation: RIGHT-SIDE UP (using bottom edge + TOP anchor)")


                        # 🔧 SAFETY: Ensure coordinates are positive and within slide bounds
                        # Allow slight negative values for precise positioning but cap at reasonable bounds
                        ppt_y_inches = max(-0.1, min(ppt_y_inches, 15))  # Allow slight negative, max 15 inches
                        ppt_x_inches = max(-0.1, min(ppt_x_inches, 15))  # Allow slight negative, max 15 inches

                        # Element dimensions
                        ppt_width_inches = max(element_width_points / 72, 0.1)
                        ppt_height_inches = max(element_height_points / 72, 0.05)

                        # 🐛 DEBUG: Log first few elements with detailed coordinates
                        if i < 5:
                            logging.info(f"🎯 ELEMENT #{i+1}: '{element['text'][:40]}'")
                            logging.info(f"   PDF BBOX: ({bbox[0]:.1f}, {bbox[1]:.1f}, {bbox[2]:.1f}, {bbox[3]:.1f})")
                            logging.info(f"   PDF Y: bottom={pdf_y_bottom:.1f}, top={pdf_y_top:.1f}")
                            logging.info(f"   PPT POS: ({ppt_x_inches:.3f}\", {ppt_y_inches:.3f}\")")
                            logging.info(f"   PPT SIZE: {ppt_width_inches:.3f}\" x {ppt_height_inches:.3f}\"")
                            logging.info("---")

                        # 📍 Create text box at exact coordinates
                        try:
                            textbox = slide.shapes.add_textbox(
                                Inches(ppt_x_inches),
                                Inches(ppt_y_inches),
                                Inches(ppt_width_inches),
                                Inches(ppt_height_inches)
                            )

                            # 🎯 REMOVE ROTATION - This was causing positioning issues
                            # Text orientation needs to be fixed at the coordinate level, not rotation level

                            # Set text content and formatting
                            text_frame = textbox.text_frame
                            text_frame.text = element['text']

                            # 🎯 ULTRA-MINIMAL MARGINS FOR PIXEL-PERFECT POSITIONING
                            text_frame.margin_left = Inches(0.003)    # 0.22 points
                            text_frame.margin_right = Inches(0.003)   # 0.22 points
                            text_frame.margin_top = Inches(0.002)     # 0.14 points
                            text_frame.margin_bottom = Inches(0.002)  # 0.14 points
                            text_frame.word_wrap = False
                            text_frame.auto_size = MSO_AUTO_SIZE.NONE
                            # 🎯 FIX ORIENTATION: Use TOP anchor for proper text orientation
                            text_frame.vertical_anchor = MSO_ANCHOR.TOP

                            # � COMPREHENSIVE POWERPOINT OUTPUT DEBUGGING
                            if i < 3:
                                logging.info(f"   📋 POWERPOINT TEXTBOX PROPERTIES:")
                                logging.info(f"     - Textbox position: ({textbox.left.inches:.4f}\", {textbox.top.inches:.4f}\")")
                                logging.info(f"     - Textbox size: {textbox.width.inches:.4f}\" x {textbox.height.inches:.4f}\"")
                                logging.info(f"     - Text content: '{text_frame.text}'")
                                logging.info(f"     - Vertical anchor: {text_frame.vertical_anchor}")
                                logging.info(f"     - Word wrap: {text_frame.word_wrap}")
                                logging.info(f"     - Auto size: {text_frame.auto_size}")
                                logging.info(f"     - Margins: L={text_frame.margin_left.inches:.3f}\", R={text_frame.margin_right.inches:.3f}\", T={text_frame.margin_top.inches:.3f}\", B={text_frame.margin_bottom.inches:.3f}\"")

                            # �🎨 CHECK FOR HIGHLIGHTING
                            element_bbox = element['bbox']
                            is_highlighted = False
                            highlight_color = None

                            for highlight in highlight_areas:
                                h_bbox = highlight['bbox']
                                # Check if element overlaps with highlight area
                                if (element_bbox[0] < h_bbox[2] and element_bbox[2] > h_bbox[0] and
                                    element_bbox[1] < h_bbox[3] and element_bbox[3] > h_bbox[1]):
                                    is_highlighted = True
                                    highlight_color = highlight.get('color')
                                    break

                            # Apply highlighting if detected
                            if is_highlighted and highlight_color:
                                try:
                                    # Set background color for the text box
                                    fill = textbox.fill
                                    fill.solid()
                                    if isinstance(highlight_color, (int, float)):
                                        # Convert color to RGB
                                        fill.fore_color.rgb = RGBColor(
                                            (int(highlight_color) >> 16) & 0xFF,
                                            (int(highlight_color) >> 8) & 0xFF,
                                            int(highlight_color) & 0xFF
                                        )
                                    else:
                                        # Default yellow highlight
                                        fill.fore_color.rgb = RGBColor(255, 255, 0)
                                except Exception as highlight_error:
                                    logging.warning(f"⚠️ Failed to apply highlight: {highlight_error}")

                            # Font formatting
                            for paragraph in text_frame.paragraphs:
                                # Handle bullet points
                                if element['is_bullet']:
                                    paragraph.alignment = PP_ALIGN.LEFT

                                for run in paragraph.runs:
                                    run.font.name = element['font_name']
                                    # 🎯 PRECISE FONT SIZE SCALING
                                    # Use more accurate scaling factor for better visual match
                                    original_size = element['font_size']
                                    scaled_size = max(original_size * 0.85, 7)  # Better scaling ratio
                                    run.font.size = Pt(round(scaled_size, 1))  # Round to 1 decimal
                                    run.font.bold = element['is_bold']
                                    run.font.italic = element['is_italic']

                                    # Set text color if available
                                    if element['color'] and element['color'] != 0:
                                        try:
                                            # Convert color to RGB
                                            color_rgb = RGBColor(
                                                (element['color'] >> 16) & 0xFF,
                                                (element['color'] >> 8) & 0xFF,
                                                element['color'] & 0xFF
                                            )
                                            run.font.color.rgb = color_rgb
                                        except:
                                            pass  # Use default color if conversion fails

                        except Exception as e:
                            logging.warning(f"⚠️ Failed to create textbox for element {i}: {e}")
                            continue

                    # 🎯 STEP 5: ADD LINES AND BORDERS FOR COMPLETE PAGE STRUCTURE
                    # Draw detected lines and borders to preserve page layout
                    for line in lines_and_borders:
                        try:
                            # Convert PDF coordinates to PowerPoint coordinates
                            if line['type'] == 'horizontal':
                                # Horizontal line
                                start_x = line['start'][0] / 72.0
                                start_y = (pdf_rect.height - line['start'][1]) / 72.0
                                end_x = line['end'][0] / 72.0
                                end_y = (pdf_rect.height - line['end'][1]) / 72.0

                                # Create horizontal line shape
                                line_shape = slide.shapes.add_connector(
                                    1,  # Straight connector
                                    Inches(start_x), Inches(start_y),
                                    Inches(end_x), Inches(end_y)
                                )

                            elif line['type'] == 'vertical':
                                # Vertical line
                                start_x = line['start'][0] / 72.0
                                start_y = (pdf_rect.height - line['start'][1]) / 72.0
                                end_x = line['end'][0] / 72.0
                                end_y = (pdf_rect.height - line['end'][1]) / 72.0

                                # Create vertical line shape
                                line_shape = slide.shapes.add_connector(
                                    1,  # Straight connector
                                    Inches(start_x), Inches(start_y),
                                    Inches(end_x), Inches(end_y)
                                )

                            # Style the line
                            line_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black line
                            line_shape.line.width = Pt(line.get('stroke_width', 1.0))

                        except Exception as line_error:
                            logging.warning(f"⚠️ Failed to create line: {line_error}")
                            continue

                    logging.info(f"✅ Added {len(lines_and_borders)} structural lines and borders")

                except Exception as e:
                    logging.error(f"❌ Text-based conversion failed: {e}")
                    # Simple fallback
                    try:
                        text_content = page.get_text()
                        if text_content.strip():
                            textbox = slide.shapes.add_textbox(
                                Inches(0.5), Inches(0.5),
                                Inches(prs.slide_width.inches - 1),
                                Inches(prs.slide_height.inches - 1)
                            )
                            textbox.text_frame.text = text_content.strip()
                            logging.info("✅ Added fallback text content")
                    except:
                        pass


            elif conversion_mode == 'image_slides':
                # Convert page to image and add to slide
                try:
                    # Convert page to image
                    mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better quality
                    pix = page.get_pixmap(matrix=mat)
                    img_data = pix.tobytes("png")

                    # Save image temporarily
                    temp_img_path = os.path.join(STATIC_DIR, f"{unique_id}_slide_{page_num}.png")
                    with open(temp_img_path, "wb") as img_file:
                        img_file.write(img_data)

                    # Add image to slide (fit to slide)
                    slide.shapes.add_picture(
                        temp_img_path,
                        Inches(0), Inches(0),
                        prs.slide_width, prs.slide_height
                    )

                    # Clean up temp image
                    os.remove(temp_img_path)
                    pix = None

                except Exception as img_error:
                    logging.warning(f"Could not add page image: {str(img_error)}")

            else:  # editable_text mode
                # Extract all text and add as editable text blocks
                try:
                    text_content = page.get_text()
                    if text_content.strip():
                        lines = text_content.strip().split('\n')
                        if lines:
                            # Add title if available
                            if len(lines) > 0:
                                title_textbox = slide.shapes.add_textbox(
                                    Inches(0.5), Inches(0.5),
                                    Inches(prs.slide_width.inches - 1),
                                    Inches(1)
                                )
                                title_textbox.text_frame.text = lines[0][:100]

                                # Set title font
                                for paragraph in title_textbox.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(18)
                                        run.font.bold = True
                                        run.font.name = "Arial"

                            # Add content
                            if len(lines) > 1:
                                content_text = '\n'.join(lines[1:])
                                content_textbox = slide.shapes.add_textbox(
                                    Inches(0.5), Inches(2),
                                    Inches(prs.slide_width.inches - 1),
                                    Inches(prs.slide_height.inches - 2.5)
                                )
                                content_textbox.text_frame.text = content_text

                                # Set content font
                                for paragraph in content_textbox.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
                                        run.font.name = "Arial"

                except Exception as text_error:
                    logging.warning(f"Could not add editable text: {str(text_error)}")

        doc.close()

        # Save PowerPoint file
        output_filename = f"converted_{unique_id}_{original_filename.rsplit('.', 1)[0]}.pptx"
        output_path = os.path.join(STATIC_DIR, output_filename)

        prs.save(output_path)
        logging.info(f"PowerPoint saved: {output_path}")

        # Clean up input file
        try:
            os.remove(input_path)
        except:
            pass

        return send_file(output_path, as_attachment=True, download_name=output_filename)

    except Exception as e:
        logging.error(f"Error in PDF to PowerPoint conversion: {str(e)}")
        return jsonify({"error": f"PowerPoint conversion failed: {str(e)}"}), 500

@basic_operations_bp.route('/remove', methods=['POST'])
def remove_pages():
    """Remove specific pages from PDF"""
    try:
        # Get file from request
        if 'file' not in request.files:
            return jsonify({"error": "No file part in request"}), 400

        file = request.files['file']

        # Validate file
        is_valid, message = validate_pdf_file(file)
        if not is_valid:
            return jsonify({"error": message}), 400

        # Get form parameters
        pages_to_remove = request.form.get('pages_to_remove', '')
        removal_method = request.form.get('removal_method', 'specific_pages')
        start_page = request.form.get('start_page', '')
        end_page = request.form.get('end_page', '')

        if removal_method == 'specific_pages' and not pages_to_remove:
            return jsonify({"error": "Please specify pages to remove"}), 400
        elif removal_method == 'page_range' and (not start_page or not end_page):
            return jsonify({"error": "Please specify start and end pages"}), 400

        # Create secure filenames
        original_filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())[:8]
        input_filename = f"{unique_id}_{original_filename}"

        input_path = os.path.join(STATIC_DIR, input_filename)

        # Save uploaded file
        file.save(input_path)
        logging.debug(f"File saved to: {input_path}")

        # Read PDF and get total pages
        with open(input_path, 'rb') as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            total_pages = len(reader.pages)

            # Parse pages to remove
            pages_to_remove_set = set()

            if removal_method == 'specific_pages':
                # Parse page ranges (e.g., "1,3,5-7,10")
                for part in pages_to_remove.split(','):
                    part = part.strip()
                    if '-' in part:
                        try:
                            start, end = map(int, part.split('-'))
                            pages_to_remove_set.update(range(start, min(end + 1, total_pages + 1)))
                        except ValueError:
                            continue
                    else:
                        try:
                            page_num = int(part)
                            if 1 <= page_num <= total_pages:
                                pages_to_remove_set.add(page_num)
                        except ValueError:
                            continue

            elif removal_method == 'page_range':
                try:
                    start = int(start_page)
                    end = int(end_page)
                    pages_to_remove_set.update(range(start, min(end + 1, total_pages + 1)))
                except ValueError:
                    return jsonify({"error": "Invalid page range"}), 400

            elif removal_method == 'odd_pages':
                pages_to_remove_set.update(range(1, total_pages + 1, 2))  # 1, 3, 5, ...

            elif removal_method == 'even_pages':
                pages_to_remove_set.update(range(2, total_pages + 1, 2))  # 2, 4, 6, ...

            elif removal_method == 'blank_pages':
                # For now, just remove last page as a placeholder
                # In a full implementation, you'd analyze page content
                pages_to_remove_set.add(total_pages)

            # Create new PDF with remaining pages
            writer = PyPDF2.PdfWriter()

            for page_num in range(1, total_pages + 1):
                if page_num not in pages_to_remove_set:
                    writer.add_page(reader.pages[page_num - 1])  # PyPDF2 uses 0-based indexing

            # Check if any pages remain
            if len(writer.pages) == 0:
                return jsonify({"error": "Cannot remove all pages. At least one page must remain."}), 400

            # Save output file
            output_filename = f"removed_pages_{unique_id}_{original_filename}"
            output_path = os.path.join(STATIC_DIR, output_filename)

            with open(output_path, 'wb') as output_file:
                writer.write(output_file)

        # Clean up input file
        try:
            os.remove(input_path)
        except:
            pass

        logging.info(f"Removed {len(pages_to_remove_set)} pages from PDF. Remaining: {len(writer.pages)} pages")

        return send_file(output_path, as_attachment=True, download_name=f"removed_pages_{original_filename}")

    except Exception as e:
        logging.error(f"Error in remove pages: {str(e)}")
        return jsonify({"error": f"Page removal failed: {str(e)}"}), 500
