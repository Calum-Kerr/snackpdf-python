from flask import Blueprint, request, send_file, jsonify
import os
import tempfile
import logging
import uuid
from datetime import datetime
from docx import Document
import openpyxl
from openpyxl.utils import get_column_letter
from pptx import Presentation
import fitz  # PyMuPDF
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.utils import ImageReader
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
import io
from PIL import Image

office_to_pdf_bp = Blueprint('office_to_pdf', __name__)

STATIC_DIR = os.path.join(os.getcwd(), 'pdf_to_image', 'static')

# Ensure the static directory exists
if not os.path.exists(STATIC_DIR):
    os.makedirs(STATIC_DIR)

# Configure logging
logging.basicConfig(level=logging.DEBUG)

# Supported file formats
WORD_FORMATS = {'.doc', '.docx'}
EXCEL_FORMATS = {'.xls', '.xlsx'}
POWERPOINT_FORMATS = {'.ppt', '.pptx'}

def validate_file(file, allowed_extensions):
    """Validate uploaded file"""
    if not file or file.filename == '':
        return False, "No file selected"
    
    file_ext = os.path.splitext(file.filename)[1].lower()
    if file_ext not in allowed_extensions:
        return False, f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}"
    
    # Check file size (max 50MB)
    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)
    
    if file_size > 50 * 1024 * 1024:  # 50MB
        return False, "File too large. Maximum size is 50MB"
    
    return True, "Valid file"

def convert_word_to_pdf(input_path, output_path):
    """Convert Word document to PDF using python-docx and reportlab"""
    try:
        # Load the Word document
        doc = Document(input_path)
        
        # Create PDF
        pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Process paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Determine style based on paragraph formatting
                if paragraph.style.name.startswith('Heading'):
                    style = styles['Heading1']
                else:
                    style = styles['Normal']
                
                p = Paragraph(paragraph.text, style)
                story.append(p)
                story.append(Spacer(1, 12))
        
        # Build PDF
        pdf_doc.build(story)
        return True, "Conversion successful"
        
    except Exception as e:
        logging.error(f"Error converting Word to PDF: {str(e)}")
        return False, f"Conversion failed: {str(e)}"

def convert_excel_to_pdf(input_path, output_path):
    """Convert Excel file to PDF"""
    try:
        # Load workbook
        workbook = openpyxl.load_workbook(input_path)
        
        # Create PDF
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        y_position = height - 50
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            
            # Add sheet title
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, y_position, f"Sheet: {sheet_name}")
            y_position -= 30
            
            # Process rows
            for row in sheet.iter_rows(values_only=True):
                if y_position < 50:  # Start new page
                    c.showPage()
                    y_position = height - 50
                
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    c.setFont("Helvetica", 10)
                    c.drawString(50, y_position, row_text[:100])  # Limit text length
                    y_position -= 15
            
            y_position -= 20  # Space between sheets
        
        c.save()
        return True, "Conversion successful"
        
    except Exception as e:
        logging.error(f"Error converting Excel to PDF: {str(e)}")
        return False, f"Conversion failed: {str(e)}"

def convert_powerpoint_to_pdf(input_path, output_path):
    """Convert PowerPoint presentation to PDF"""
    try:
        # Load presentation
        prs = Presentation(input_path)
        
        # Create PDF
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        for i, slide in enumerate(prs.slides):
            if i > 0:
                c.showPage()
            
            # Add slide number
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, height - 50, f"Slide {i + 1}")
            
            y_position = height - 100
            
            # Process text shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    c.setFont("Helvetica", 12)
                    # Split long text into multiple lines
                    text_lines = shape.text.split('\n')
                    for line in text_lines:
                        if y_position < 50:
                            break
                        c.drawString(50, y_position, line[:80])  # Limit line length
                        y_position -= 20
                    y_position -= 10
        
        c.save()
        return True, "Conversion successful"
        
    except Exception as e:
        logging.error(f"Error converting PowerPoint to PDF: {str(e)}")
        return False, f"Conversion failed: {str(e)}"

@office_to_pdf_bp.route('/api/word_to_pdf', methods=['POST'])
def word_to_pdf():
    """Convert Word document to PDF"""
    try:
        # Validate file
        file = request.files.get('file')
        is_valid, message = validate_file(file, WORD_FORMATS)
        if not is_valid:
            return jsonify({"error": message}), 400
        
        # Generate unique filename
        unique_id = str(uuid.uuid4())[:8]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"converted_word_{timestamp}_{unique_id}.pdf"
        output_path = os.path.join(STATIC_DIR, output_filename)
        
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1]) as temp_file:
            file.save(temp_file.name)
            
            # Convert to PDF
            success, message = convert_word_to_pdf(temp_file.name, output_path)
            
            # Clean up temp file
            os.unlink(temp_file.name)
            
            if not success:
                return jsonify({"error": message}), 500
        
        return send_file(output_path, as_attachment=True, download_name=output_filename)
        
    except Exception as e:
        logging.error(f"Error in word_to_pdf: {str(e)}")
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

@office_to_pdf_bp.route('/api/excel_to_pdf', methods=['POST'])
def excel_to_pdf():
    """Convert Excel file to PDF"""
    try:
        # Validate file
        file = request.files.get('file')
        is_valid, message = validate_file(file, EXCEL_FORMATS)
        if not is_valid:
            return jsonify({"error": message}), 400
        
        # Generate unique filename
        unique_id = str(uuid.uuid4())[:8]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"converted_excel_{timestamp}_{unique_id}.pdf"
        output_path = os.path.join(STATIC_DIR, output_filename)
        
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1]) as temp_file:
            file.save(temp_file.name)
            
            # Convert to PDF
            success, message = convert_excel_to_pdf(temp_file.name, output_path)
            
            # Clean up temp file
            os.unlink(temp_file.name)
            
            if not success:
                return jsonify({"error": message}), 500
        
        return send_file(output_path, as_attachment=True, download_name=output_filename)
        
    except Exception as e:
        logging.error(f"Error in excel_to_pdf: {str(e)}")
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

@office_to_pdf_bp.route('/api/powerpoint_to_pdf', methods=['POST'])
def powerpoint_to_pdf():
    """Convert PowerPoint presentation to PDF"""
    try:
        # Validate file
        file = request.files.get('file')
        is_valid, message = validate_file(file, POWERPOINT_FORMATS)
        if not is_valid:
            return jsonify({"error": message}), 400
        
        # Generate unique filename
        unique_id = str(uuid.uuid4())[:8]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"converted_powerpoint_{timestamp}_{unique_id}.pdf"
        output_path = os.path.join(STATIC_DIR, output_filename)
        
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1]) as temp_file:
            file.save(temp_file.name)
            
            # Convert to PDF
            success, message = convert_powerpoint_to_pdf(temp_file.name, output_path)
            
            # Clean up temp file
            os.unlink(temp_file.name)
            
            if not success:
                return jsonify({"error": message}), 500
        
        return send_file(output_path, as_attachment=True, download_name=output_filename)
        
    except Exception as e:
        logging.error(f"Error in powerpoint_to_pdf: {str(e)}")
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500
