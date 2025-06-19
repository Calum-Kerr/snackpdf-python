from flask import Blueprint, request, send_file, jsonify
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os
import tempfile
import logging
import shutil

pdf_to_powerpoint_bp = Blueprint('pdf_to_powerpoint', __name__)

STATIC_DIR = os.path.join(os.getcwd(), 'pdf_to_image', 'static')

# Ensure the static directory exists
if not os.path.exists(STATIC_DIR):
    os.makedirs(STATIC_DIR)

def convert_pdf_to_powerpoint(pdf_path, pptx_path, options=None):
    """Convert PDF to PowerPoint presentation"""
    if options is None:
        options = {}

    # Get options with defaults
    dpi = int(options.get('dpi', 300))
    slide_layout = options.get('slide_layout', 'blank')  # blank, title_only, title_content
    maintain_aspect_ratio = options.get('maintain_aspect_ratio', True)
    image_quality = options.get('image_quality', 'high')  # high, medium, low

    try:
        # Open PDF
        pdf_doc = fitz.open(pdf_path)

        # Create PowerPoint presentation
        prs = Presentation()

        # Set zoom level based on DPI
        zoom = dpi / 72.0

        for page_num in range(len(pdf_doc)):
            page = pdf_doc.load_page(page_num)

            # Convert page to image
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")

            # Create temporary image file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_img:
                temp_img.write(img_data)
                temp_img_path = temp_img.name

            try:
                # Add slide to presentation
                if slide_layout == 'blank':
                    slide_layout_obj = prs.slide_layouts[6]  # Blank layout
                elif slide_layout == 'title_only':
                    slide_layout_obj = prs.slide_layouts[5]  # Title only layout
                else:
                    slide_layout_obj = prs.slide_layouts[6]  # Default to blank

                slide = prs.slides.add_slide(slide_layout_obj)

                # Calculate image dimensions for slide
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                # Open image to get dimensions
                with Image.open(temp_img_path) as img:
                    img_width, img_height = img.size

                if maintain_aspect_ratio:
                    # Calculate scaling to fit slide while maintaining aspect ratio
                    width_ratio = slide_width / img_width
                    height_ratio = slide_height / img_height
                    scale_ratio = min(width_ratio, height_ratio)

                    new_width = int(img_width * scale_ratio)
                    new_height = int(img_height * scale_ratio)

                    # Center the image on the slide
                    left = (slide_width - new_width) // 2
                    top = (slide_height - new_height) // 2
                else:
                    # Stretch to fill entire slide
                    new_width = slide_width
                    new_height = slide_height
                    left = 0
                    top = 0

                # Add image to slide
                slide.shapes.add_picture(temp_img_path, left, top, new_width, new_height)

                # Add title if using title layout
                if slide_layout == 'title_only' and hasattr(slide.shapes, 'title'):
                    slide.shapes.title.text = f"Page {page_num + 1}"

            finally:
                # Clean up temporary image file
                if os.path.exists(temp_img_path):
                    os.unlink(temp_img_path)

        # Save PowerPoint presentation
        prs.save(pptx_path)
        pdf_doc.close()

        return True

    except Exception as e:
        logging.error(f"Error converting PDF to PowerPoint: {e}")
        raise e

@pdf_to_powerpoint_bp.route('/pdf_to_powerpoint', methods=['POST'])
def pdf_to_powerpoint():
    """Handle PDF to PowerPoint conversion requests"""
    try:
        file = request.files['file']

        # Get options from form
        options = {
            'dpi': int(request.form.get('dpi', 300)),
            'slide_layout': request.form.get('slide_layout', 'blank'),
            'maintain_aspect_ratio': request.form.get('maintain_aspect_ratio', 'true').lower() == 'true',
            'image_quality': request.form.get('image_quality', 'high')
        }

        output_dir = os.path.join(STATIC_DIR, 'powerpoint')

        if os.path.exists(output_dir):
            shutil.rmtree(output_dir)
        os.makedirs(output_dir, exist_ok=True)

        with tempfile.TemporaryDirectory() as temp_dir:
            # Save uploaded PDF
            pdf_path = os.path.join(temp_dir, file.filename)
            file.save(pdf_path)

            # Create output PowerPoint path
            base_name = os.path.splitext(file.filename)[0]
            pptx_path = os.path.join(output_dir, f'{base_name}.pptx')

            # Convert PDF to PowerPoint
            convert_pdf_to_powerpoint(pdf_path, pptx_path, options)

            logging.debug(f"PDF converted to PowerPoint: {pptx_path}")

            return send_file(pptx_path, as_attachment=True,
                           download_name=f'{base_name}.pptx',
                           mimetype='application/vnd.openxmlformats-presentationml.presentation')

    except Exception as e:
        logging.error(f"Error in PDF to PowerPoint conversion: {e}")
        return jsonify({"error": str(e)}), 500